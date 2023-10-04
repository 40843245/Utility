import pptx

from StringHandler import StringHandler

class PowerPointHandler():
    def __init__(self):
        pass
    def SetInputFullName(self,inputFullName : str):
        self.inputFullName = inputFullName
    def Load(self):
        self.powerpoint = pptx.Presentation(self.inputFullName)
    def GetSlider(self,pageIndex):
        r = self.GetSliders()
        retVal = [ (page,slide) for (page,slide) in r if page == pageIndex ]
        return retVal
    def GetSliders(self):
        notes = list()
        for page, slide in enumerate(self.powerpoint.slides):
            temp = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    temp.append(shape.text)
            notes.append((page,temp))    
        return notes
    def SearchAll(self,searchText : str):
        delimList = [ ' ','\n','\t']
        result = list()
        notes = self.GetSliders()
        for i in range(0,len(notes),1):
            pagePair = notes[i]
            contents = pagePair[1]
            for content in contents:
                r = StringHandler.Split(content, delimList)
                for j in range(0,len(r),1):
                    if r[j][2] == searchText:
                        x = ( i , ) + r[j] 
                        result.append( x  )
        return result
                
    def ReplaceAll(self,searchText:str,replaceText:str):
        delimList = [ ' ','\n','\t']
        resultList3 = list()
        notes = self.GetSliders()
        for i in range(0,len(notes),1):
            resultList2 = list()
            pagePair = notes[i]
            contents = pagePair[1]
            for content in contents:
                resultList1 = list()
                r = StringHandler.Split(content, delimList)
                for j in range(0,len(r),1):
                    if r[j][2] == searchText:
                       s = replaceText 
                       resultList1.append(s) 
                    else:
                        s = r [j][2]
                        resultList1.append(s) 
                    if j != len(r)-1:
                        resultList1.append(delimList[r[j][1]]) 
                        
                    x = ''.join(resultList1)
                resultList2.append(x)
            resultList3.append(resultList2)
        return resultList3
        
         
if __name__ == '__main__':
    inputFullName = r"C:\QtProject\Test2\test1\test.pptx"
    
    powerPowerPointHandler = PowerPointHandler()
    powerPowerPointHandler.SetInputFullName(inputFullName)
    powerPowerPointHandler.Load()
    
    notes = powerPowerPointHandler.GetSliders()
    print(notes)
    
    r = powerPowerPointHandler.SearchAll("Adding")
    print(r)
    
    r = powerPowerPointHandler.ReplaceAll("Adding","Eating")
    print(r)
    
    
    
        
